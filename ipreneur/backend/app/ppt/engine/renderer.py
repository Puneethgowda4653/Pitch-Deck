"""
PPT Renderer v2 — Gamma-style visual presentation engine.

Layout system:
  full_bleed  — Cover/closing: centered title, gradient bar, minimal text
  big_number  — Market/traction: 3-4 large metric callouts in bordered cards
  title_bullets — Standard content: title + lead-in bullets in two columns
  two_column  — Comparison: left vs right with heading + bullet list
  cards       — Feature/model: 2-4 styled cards with title, body, optional metric
  timeline    — Milestones: horizontal flow of events
"""
import io
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from loguru import logger

from app.services.branding.extractor import BrandingData
from app.agents.content.content_agent import DeckContent, SlideContent
from app.ppt.engine.themes import get_theme
from app.ppt.engine.template_schema import TEMPLATE_SECTION_BY_SLIDE_NUMBER, _coerce_number as _coerce_float
from app.services.images.image_service import normalize_logo


# ── Color helpers ──────────────────────────────────────────────────────────────

def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _tint(c: RGBColor, factor: float) -> RGBColor:
    """Darken color toward black by factor (0 = black, 1 = original)."""
    return RGBColor(int(c[0] * factor), int(c[1] * factor), int(c[2] * factor))


def _mix(c1: RGBColor, c2: RGBColor, t: float) -> RGBColor:
    """Linear interpolation between two colors (t=0 → c1, t=1 → c2)."""
    return RGBColor(
        int(c1[0] * (1 - t) + c2[0] * t),
        int(c1[1] * (1 - t) + c2[1] * t),
        int(c1[2] * (1 - t) + c2[2] * t),
    )


def _brightness(c: RGBColor) -> float:
    return (c[0] * 299 + c[1] * 587 + c[2] * 114) / 1000


def _ensure_contrast(fg: RGBColor, bg: RGBColor, min_diff: float = 90.0) -> RGBColor:
    """Nudge `fg` toward black or white (whichever the background needs) until
    it's readably distinct from `bg` — works for both light and dark themes,
    unlike a one-directional "always brighten" check."""
    diff = abs(_brightness(fg) - _brightness(bg))
    if diff >= min_diff:
        return fg
    target = _rgb("#FFFFFF") if _brightness(bg) < 128 else _rgb("#000000")
    # Blend toward the target just enough to clear the contrast floor.
    t = min(1.0, (min_diff - diff) / 255 + 0.15)
    return _mix(fg, target, t)


_DEFAULT_PRIMARY, _DEFAULT_SECONDARY = "#2540B5", "#7B2CBF"  # BrandingResult fallback constants


def _apply_brand_accents(palette: dict, branding: Optional["BrandingData"]) -> dict:
    """Override the theme's accent colors with the company's real brand colors,
    keeping the theme's bg/surface/text tokens (and light/dark correctness)
    untouched. Skips the override when branding colors are still the generic
    hardcoded defaults — i.e. no real brand signal was ever found for this
    company, so the curated theme palette is a better bet than a fake accent."""
    if not branding:
        return palette
    pc = getattr(branding, "primary_color", None)
    sc = getattr(branding, "secondary_color", None)
    if not pc or not sc or (pc == _DEFAULT_PRIMARY and sc == _DEFAULT_SECONDARY):
        return palette
    try:
        primary = _ensure_contrast(_rgb(pc), palette["bg"])
        secondary = _ensure_contrast(_rgb(sc), palette["bg"])
    except Exception:
        return palette
    out = dict(palette)
    out["accent1"] = primary
    out["accent2"] = secondary
    out["accent3"] = _mix(primary, secondary, 0.5)
    out["accent4"] = _mix(primary, palette["text_pri"], 0.25)
    return out


def _set_fill_alpha(shape, color: RGBColor, alpha_pct: float) -> None:
    """Set real alpha transparency on a shape's solid fill via direct XML —
    python-pptx's high-level API doesn't expose this, but the OOXML element
    supports it directly. `alpha_pct` is 0-100."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid_fill = shape._element.spPr.find(f"{{{ns}}}solidFill")
    srgb = solid_fill.find(f"{{{ns}}}srgbClr") if solid_fill is not None else None
    if srgb is None:
        return
    from lxml import etree
    etree.SubElement(srgb, f"{{{ns}}}alpha").set("val", str(int(max(0, min(100, alpha_pct)) * 1000)))


_MONEY_RE = None  # compiled lazily below to keep module import cheap


def _parse_money_to_millions(raw) -> Optional[float]:
    """'$47B' -> 47000.0, '$640M' -> 640.0, '1.2M' -> 1.2 (bare/'M' assumed
    millions, consistent with traction.series already being documented as $M).
    Used for chart bar length — the original string stays the visible label."""
    global _MONEY_RE
    if not isinstance(raw, str):
        return None
    if _MONEY_RE is None:
        import re
        _MONEY_RE = re.compile(r"[-+]?[\d,]*\.?\d+")
    s = raw.strip().upper().replace(",", "").replace("$", "")
    m = _MONEY_RE.search(s)
    if not m:
        return None
    value = float(m.group())
    if "B" in s:
        value *= 1000
    elif "K" in s:
        value /= 1000
    elif "T" in s:
        value *= 1_000_000
    return value


def _theme_palette(theme: dict) -> dict:
    """
    Build the slide palette from a selected deck theme (see themes.py).

    Uses the theme's own tokens directly so the PPTX matches the template the
    user previewed in the module — including light themes (light background +
    dark text), instead of the old always-dark brand palette.
    """
    bg      = _rgb(theme["bg"])
    surface = _rgb(theme["surface"])
    text    = _rgb(theme["text"])
    muted   = _rgb(theme["muted"])
    accent  = _rgb(theme["accent"])
    accent2 = _rgb(theme["accent2"])

    return {
        "mode":      theme.get("mode", "dark"),
        "bg":        bg,
        "surface":   surface,
        # Alt surface — nudge slightly toward the text color for contrast
        "surface2":  _mix(surface, text, 0.06),
        # Border — a soft blend of surface toward text (works in light & dark)
        "border":    _mix(surface, text, 0.20),
        "text_pri":  text,
        "text_sec":  muted,
        # Muted text — faded toward the background
        "text_muted": _mix(muted, bg, 0.40),
        "accent1":   accent,
        "accent2":   accent2,
        "accent3":   _mix(accent, accent2, 0.5),
        "accent4":   _mix(accent, text, 0.25),
        "font_title": theme.get("font_title", "Calibri"),
        "font_body":  theme.get("font_body", "Calibri"),
        "title_upper": bool(theme.get("title_upper", False)),
    }


# ── Constants (fallback defaults) ─────────────────────────────────────────────

W = Inches(13.333)
H = Inches(7.5)

TEXT_PRI    = _rgb("#F5F3FF")
TEXT_MUTED  = _rgb("#5C4E7A")
ACCENT_1    = _rgb("#2540B5")
ACCENT_2    = _rgb("#7B2CBF")


class PPTRenderer:
    """Renders a DeckContent into a PPTX file using branded, Gamma-style layouts."""

    def __init__(self):
        pass

    async def render(
        self,
        deck_content: DeckContent,
        branding: BrandingData,
        logo_bytes: Optional[bytes] = None,
        slide_backgrounds: Optional[dict] = None,
        template_key: Optional[str] = None,
        template_data: Optional[dict] = None,
    ) -> bytes:
        theme = get_theme(template_key)
        logger.info(
            f"📊 Rendering PPTX: {len(deck_content.slides)} slides | "
            f"theme={template_key or 'default'} ({theme['name']}, {theme['mode']}) | "
            f"template_data={'yes' if template_data else 'no'}"
        )

        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H

        # Palette comes from the user-selected theme so the download matches the
        # template previewed in the module (dark OR light) — then the company's
        # real brand colors override the accent tokens, if any were found.
        palette = _apply_brand_accents(_theme_palette(theme), branding)
        slide_backgrounds = slide_backgrounds or {}
        template_data = template_data or {}
        # Normalize once here so every placement site downstream can assume
        # PNG bytes — python-pptx can't embed SVG, which previously failed
        # silently at each add_picture() call.
        if logo_bytes:
            logo_bytes = normalize_logo(logo_bytes)

        for slide_data in deck_content.slides:
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            _bg(slide, palette["bg"])

            # Add background photo if available for this slide type
            bg_bytes = slide_backgrounds.get(slide_data.slide_type)
            if bg_bytes:
                self._add_background_image(slide, bg_bytes, palette, overlay_alpha_pct=50)

            layout = slide_data.layout or "title_bullets"
            # Richer AI output for this slide, keyed by fixed slide position
            # (see template_schema.py) — each builder falls back to plain
            # SlideContent fields when this is empty/missing for its slide.
            section_key = TEMPLATE_SECTION_BY_SLIDE_NUMBER.get(slide_data.slide_number)
            td = template_data.get(section_key, {}) if section_key else {}

            if layout == "full_bleed":
                self._full_bleed(slide, slide_data, palette, logo_bytes=logo_bytes if slide_data.slide_number == 1 else None, td=template_data)
            elif layout == "market_sizing" or slide_data.slide_type == "financials":
                self._market_sizing(slide, slide_data, palette, td=td)
            elif layout == "big_number":
                self._big_number(slide, slide_data, palette, td=td)
            elif layout == "cards":
                self._cards(slide, slide_data, palette, td=td)
            elif layout == "two_column":
                self._two_column(slide, slide_data, palette, td=td)
            elif layout == "timeline":
                self._timeline(slide, slide_data, palette)
            else:
                self._title_bullets(slide, slide_data, palette)

            # Add small logo to top-right of every non-cover slide
            if logo_bytes and slide_data.slide_number != 1:
                self._add_logo_small(slide, logo_bytes)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        data = buf.read()
        logger.info(f"✅ PPTX rendered: {len(data):,} bytes | logo={'yes' if logo_bytes else 'no'} | bg_photos={len(slide_backgrounds)}")
        return data

    # ── Layout renderers ───────────────────────────────────────────────────────

    def _full_bleed(self, slide, s: SlideContent, p: dict, logo_bytes: Optional[bytes] = None, td: Optional[dict] = None):
        """Full-screen cover/closing slide — with optional logo."""
        td = td or {}
        title_font = p["font_title"]
        _rect(slide, 0, 0, W, Inches(0.12), p["accent1"])
        _rect(slide, 0, H - Inches(0.06), W, Inches(0.06), p["accent2"])

        # Logo on cover — top-center
        if logo_bytes:
            try:
                logo_stream = io.BytesIO(logo_bytes)
                native = _image_size(logo_bytes)
                logo_w, logo_h = _fit_box(native, Inches(2.4), Inches(0.9)) if native else (Inches(2.0), Inches(0.8))
                logo_x = (W - logo_w) / 2
                slide.shapes.add_picture(logo_stream, logo_x, Inches(0.3), logo_w, logo_h)
                title_y = Inches(1.4)
            except Exception:
                title_y = Inches(1.8)
        else:
            title_y = Inches(1.8)

        title_text = s.title.upper() if p.get("title_upper") else s.title
        _txt(slide, Inches(1), title_y, Inches(11.3), Inches(1.6),
             title_text, size=48, bold=True, color=p["text_pri"], align=PP_ALIGN.CENTER, font=title_font)

        # Cover-only enrichment from template_data's top-level identity fields
        # (round/tagline) — only on slide 1, and only if SlideContent didn't
        # already carry a subtitle (never overrides real AI-generated content).
        subtitle = s.subtitle or (td.get("round") or td.get("tagline") if s.slide_number == 1 else "")
        if subtitle:
            _txt(slide, Inches(1.5), title_y + Inches(1.7), Inches(10.3), Inches(0.8),
                 subtitle, size=22, color=p["text_sec"], align=PP_ALIGN.CENTER)

        if s.body:
            _txt(slide, Inches(2), title_y + Inches(2.6), Inches(9.3), Inches(0.8),
                 s.body, size=15, color=p["text_muted"], align=PP_ALIGN.CENTER)

        closing = td.get("closing") or {}
        footer = (closing.get("contact") if s.slide_number != 1 and closing.get("contact") else None)
        _txt(slide, Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.4),
             footer or "Investor Presentation  •  Confidential",
             size=10, color=p["text_muted"], align=PP_ALIGN.CENTER)

    def _big_number(self, slide, s: SlideContent, p: dict, td: Optional[dict] = None):
        """Market/traction/summary slide with large metric callouts, and — when
        template_data supplies a real numeric series (the Traction slide) — a
        native, editable column chart instead of a bullet-point list."""
        td = td or {}
        self._slide_header(slide, s, p)

        # Prefer template_data's KPI tiles (summary.highlights / traction.kpis,
        # both {k,l} shape) over the flatter SlideContent.data_points.
        td_kpis = td.get("kpis") or td.get("highlights") or []
        if td_kpis:
            metrics = [{"value": k.get("k", ""), "label": k.get("l", "")} for k in td_kpis[:4]]
        else:
            metrics = s.data_points[:4] if s.data_points else []
        n = len(metrics)

        if n > 0:
            card_w = Inches(2.8)
            total = n * 2.8 + (n - 1) * 0.25
            start_x = (13.333 - total) / 2
            accents = [p["accent1"], p["accent2"], p["accent3"], p["accent4"]]

            for i, dp in enumerate(metrics):
                x = Inches(start_x + i * (2.8 + 0.25))
                y = Inches(2.3)
                ch = Inches(2.4)

                card = _rect(slide, x, y, card_w, ch, p["surface"])
                card.line.color.rgb = accents[i % len(accents)]
                card.line.width = Pt(1.5)

                _txt(slide, x + Inches(0.15), y + Inches(0.25), Inches(2.5), Inches(1.0),
                     dp.get("value", ""), size=34, bold=True, color=p["text_pri"], align=PP_ALIGN.CENTER)

                _txt(slide, x + Inches(0.15), y + Inches(1.25), Inches(2.5), Inches(0.45),
                     dp.get("label", ""), size=12, color=p["text_sec"], align=PP_ALIGN.CENTER)

                if dp.get("sublabel"):
                    _txt(slide, x + Inches(0.1), y + Inches(1.7), Inches(2.6), Inches(0.5),
                         dp["sublabel"], size=10, color=p["text_muted"], align=PP_ALIGN.CENTER)

        series = td.get("series") or []
        chart_drawn = False
        if len(series) >= 2:
            chart_drawn = self._traction_chart(
                slide, series, p, x=Inches(0.8), y=Inches(4.9), cx=Inches(11.7), cy=Inches(1.9)
            )

        if not chart_drawn:
            if s.body:
                _txt(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.9),
                     s.body, size=14, color=p["text_sec"])

            if s.bullet_points:
                y_b = Inches(5.9) if s.body else Inches(5.2)
                for i, bp in enumerate(s.bullet_points[:3]):
                    _txt(slide, Inches(1.0), y_b + Inches(i * 0.45), Inches(11.3), Inches(0.4),
                         f"›  {bp}", size=13, color=p["text_sec"])

    def _traction_chart(self, slide, series: list, p: dict, x, y, cx, cy) -> bool:
        """Native, editable column chart for a quarterly growth series
        (template_data.traction.series — discrete snapshots, so columns are a
        more honest representation than an interpolated line)."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        cats = [str(pt.get("y", "")) for pt in series]
        vals = [_coerce_float(pt.get("v")) for pt in series]
        if len(vals) < 2 or all(v == 0 for v in vals):
            return False

        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("Growth ($M)", vals)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, cd).chart
        chart.has_legend = False
        chart.value_axis.visible = False
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.category_axis.tick_labels.font.color.rgb = p["text_sec"]
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(11)
        plot.data_labels.font.color.rgb = p["text_pri"]
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = p["accent1"]
        return True

    def _cards(self, slide, s: SlideContent, p: dict, td: Optional[dict] = None):
        """Feature/model/problem/team/competition/ask slide — dispatches to a
        real comparison table (Competitive Landscape) or doughnut chart (Ask)
        when template_data supplies that shape, else a 2-4 card grid sourced
        from the richest content available (team bios / product steps /
        business-model streams+tiers / plain SlideContent.cards)."""
        td = td or {}
        self._slide_header(slide, s, p)

        if td.get("rows") and td.get("cols"):
            if self._competition_table(slide, td, p, x=Inches(0.6), y=Inches(2.0), cx=Inches(12.1), cy=Inches(4.6)):
                return

        if td.get("use"):
            if self._ask_chart(slide, td, p, x=Inches(0.8), y=Inches(2.0), cx=Inches(5.5), cy=Inches(4.6)):
                return

        cards = self._resolve_cards(s, td)
        n = len(cards)
        if n == 0:
            self._title_bullets(slide, s, p)
            return

        if n <= 2:
            card_w, card_h = Inches(5.6), Inches(3.8)
            gap = Inches(0.4)
            total_w = n * 5.6 + (n - 1) * 0.4
            y = Inches(2.0)
        else:
            card_w, card_h = Inches(3.0), Inches(3.6)
            gap = Inches(0.25)
            total_w = n * 3.0 + (n - 1) * 0.25
            y = Inches(2.1)

        start_x = (13.333 - total_w) / 2
        accents = [p["accent1"], p["accent2"], p["accent3"], p["accent4"]]

        for i, card in enumerate(cards):
            x = Inches(start_x + i * (card_w.inches + gap.inches))
            acc = accents[i % len(accents)]

            card_bg = _rect(slide, x, y, card_w, card_h, p["surface"])
            card_bg.line.color.rgb = p["border"]
            card_bg.line.width = Pt(1)

            _rect(slide, x, y, card_w, Inches(0.06), acc)

            _txt(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.6),
                 card.get("title", ""), size=16, bold=True, color=p["text_pri"], font=p["font_title"])

            metric = card.get("metric", "")
            if metric:
                _txt(slide, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.55),
                     metric, size=22, bold=True, color=acc)
                body_y = y + Inches(1.35)
            else:
                body_y = y + Inches(0.85)

            if card.get("body"):
                _txt(slide, x + Inches(0.2), body_y, card_w - Inches(0.4),
                     card_h - (body_y - y) - Inches(0.2),
                     card["body"], size=13, color=p["text_sec"], wrap=True)

    @staticmethod
    def _resolve_cards(s: SlideContent, td: dict) -> list:
        """Prefer richer template_data content (team bios / product steps /
        business-model streams+tiers) over the flatter SlideContent.cards,
        normalized into the {title, metric, body} shape the card grid uses."""
        if td.get("members"):
            return [
                {
                    "title": f"{m.get('n', '')} — {m.get('r', '')}".strip(" —"),
                    "metric": m.get("i", ""),
                    "body": m.get("b", ""),
                }
                for m in td["members"][:4]
            ]
        if td.get("steps"):
            return [
                {
                    "title": f"{st.get('n', '')}. {st.get('t', '')}".strip(". "),
                    "metric": "",
                    "body": " · ".join([st.get("d", ""), *st.get("tags", [])]).strip(" ·"),
                }
                for st in td["steps"][:4]
            ]
        if td.get("streams") or td.get("tiers"):
            cards = [
                {
                    "title": t.get("t", ""),
                    "metric": t.get("v", ""),
                    "body": (t.get("d", "") + (f" · {t.get('vl', '')}" if t.get("vl") else "")).strip(),
                }
                for t in (td.get("streams") or [])[:2]
            ]
            cards += [
                {"title": t.get("t", ""), "metric": f"{t.get('p', '')}{t.get('s', '')}", "body": t.get("d", "")}
                for t in (td.get("tiers") or [])[:2]
            ]
            if cards:
                return cards[:4]
        return s.cards[:4] if s.cards else []

    def _competition_table(self, slide, td: dict, p: dict, x, y, cx, cy) -> bool:
        """Real comparison table for template_data.competition — this is
        structurally a table (company columns × feature-boolean rows), not
        cards. Emphasizes column 0 ("This company") with bold accent checks."""
        cols = [str(c) for c in (td.get("cols") or [])][:4]
        rows = [r for r in (td.get("rows") or []) if isinstance(r, dict)][:6]
        if len(cols) < 2 or len(rows) < 2:
            return False

        n_rows, n_cols = len(rows) + 1, len(cols) + 1
        table = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy).table
        table.columns[0].width = Emu(int(cx * 0.34))
        remaining = cx - table.columns[0].width
        for ci in range(1, n_cols):
            table.columns[ci].width = Emu(int(remaining / (n_cols - 1)))

        header_text_color = _rgb("#FFFFFF") if _brightness(p["accent1"]) < 140 else _rgb("#000000")
        table.cell(0, 0).text = ""
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = p["accent1"]
        for ci, name in enumerate(cols):
            cell = table.cell(0, ci + 1)
            cell.text = name
            cell.fill.solid()
            cell.fill.fore_color.rgb = p["accent1"]
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = header_text_color

        for ri, row in enumerate(rows):
            row_bg = p["surface"] if ri % 2 == 0 else p["surface2"]
            fcell = table.cell(ri + 1, 0)
            fcell.text = str(row.get("f", ""))
            fcell.fill.solid()
            fcell.fill.fore_color.rgb = row_bg
            fpara = fcell.text_frame.paragraphs[0]
            fpara.font.size = Pt(12)
            fpara.font.color.rgb = p["text_pri"]

            flags = row.get("v") or []
            for ci in range(len(cols)):
                cell = table.cell(ri + 1, ci + 1)
                is_first_col = ci == 0
                has_it = bool(flags[ci]) if ci < len(flags) else False
                cell.text = "✓" if has_it else "—"
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                cpara = cell.text_frame.paragraphs[0]
                cpara.alignment = PP_ALIGN.CENTER
                cpara.font.size = Pt(13)
                cpara.font.bold = has_it and is_first_col
                cpara.font.color.rgb = (p["accent1"] if has_it and is_first_col
                                         else p["text_sec"] if has_it else p["text_muted"])
        return True

    def _ask_chart(self, slide, td: dict, p: dict, x, y, cx, cy) -> bool:
        """Native doughnut chart for fund allocation (template_data.ask.use) —
        the one slide where a doughnut is semantically correct, since these
        percentages genuinely sum to 100%, unlike TAM/SAM/SOM."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        items = [u for u in (td.get("use") or []) if u.get("l")][:4]
        if len(items) < 2:
            return False

        cd = CategoryChartData()
        cd.categories = [u["l"] for u in items]
        cd.add_series("Use of funds", [_coerce_float(u.get("p")) for u in items])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, cd).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.color.rgb = p["text_sec"]
        chart.legend.font.size = Pt(12)
        accents = [p["accent1"], p["accent2"], p["accent3"], p["accent4"]]
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = '0"%"'
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.color.rgb = p["text_pri"]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = accents[i % len(accents)]

        headline = td.get("headline", "")
        if headline:
            _txt(slide, x + cx + Inches(0.3), y + Inches(0.3), Inches(4.5), Inches(1.0),
                 headline, size=20, bold=True, color=p["text_pri"], font=p["font_title"])
        return True

    def _two_column(self, slide, s: SlideContent, p: dict, td: Optional[dict] = None):
        """Side-by-side comparison layout (Problem & Solution)."""
        td = td or {}
        self._slide_header(slide, s, p)

        # Prefer template_data.probsol (bold-lead {k,t} items + a footer line
        # per column) over the flatter SlideContent.columns.
        if td.get("problem") and td.get("solution"):
            cols = [
                {
                    "heading": td.get("problemTitle") or "The Problem",
                    "highlight": td.get("problemLead", ""),
                    "points": [f"**{i.get('k', '')}:** {i.get('t', '')}" for i in td["problem"][:4]],
                    "foot": td.get("problemFoot", ""),
                },
                {
                    "heading": td.get("solutionTitle") or "Our Solution",
                    "highlight": td.get("solutionLead", ""),
                    "points": [f"**{i.get('k', '')}:** {i.get('t', '')}" for i in td["solution"][:4]],
                    "foot": td.get("solutionFoot", ""),
                },
            ]
        else:
            cols = s.columns[:2] if s.columns else []
        if not cols:
            self._title_bullets(slide, s, p)
            return

        col_w, col_h = Inches(5.8), Inches(4.5)
        y = Inches(1.9)
        left_x = Inches(0.7)
        right_x = left_x + col_w + Inches(0.5)
        accents = [p["accent1"], p["accent2"]]

        for ci, (col, x) in enumerate([(cols[0], left_x), (cols[1] if len(cols) > 1 else {}, right_x)]):
            acc = accents[ci]

            bg = _rect(slide, x, y, col_w, col_h, p["surface"])
            bg.line.color.rgb = acc
            bg.line.width = Pt(1.5)

            _rect(slide, x, y, col_w, Inches(0.07), acc)

            _txt(slide, x + Inches(0.2), y + Inches(0.2), col_w - Inches(0.4), Inches(0.55),
                 col.get("heading", ""), size=18, bold=True, color=p["text_pri"], font=p["font_title"])

            yp = y + Inches(0.8)
            if col.get("highlight"):
                _txt(slide, x + Inches(0.2), yp, col_w - Inches(0.4), Inches(0.55),
                     col["highlight"], size=24, bold=True, color=acc)
                yp += Inches(0.6)

            for pt in (col.get("points") or [])[:5]:
                _txt(slide, x + Inches(0.3), yp, col_w - Inches(0.5), Inches(0.48),
                     f"•  {pt}", size=13, color=p["text_sec"])
                yp += Inches(0.48)

            if col.get("foot"):
                _txt(slide, x + Inches(0.2), y + col_h - Inches(0.5), col_w - Inches(0.4), Inches(0.4),
                     col["foot"], size=11, color=p["text_muted"])

        div = slide.shapes.add_shape(1,
            left_x + col_w + Inches(0.22), y + Inches(0.3),
            Inches(0.03), col_h - Inches(0.6))
        div.fill.solid()
        div.fill.fore_color.rgb = p["border"]
        div.line.fill.background()

    def _market_sizing(self, slide, s: SlideContent, p: dict, td: Optional[dict] = None):
        """TAM / SAM / SOM slide — a real, correctly-scaled bar chart (not a
        pie/doughnut: TAM/SAM/SOM are nested subsets, not sibling slices
        summing to a whole) + metric legend."""
        td = td or {}
        self._slide_header(slide, s, p)

        metrics = (s.data_points or [])[:4]
        labels = ["TAM", "SAM", "SOM"]
        ring_colors = [p["accent1"], p["accent2"], p["accent3"]]

        # Prefer template_data.market (tam/sam/som, each {v,l}) for both the
        # chart and legend; fall back to SlideContent.data_points for either
        # piece independently if template_data is missing/insufficient.
        td_points = []
        for key in ("tam", "sam", "som"):
            point = td.get(key) or {}
            if point.get("v"):
                td_points.append({"value": point["v"], "sublabel": point.get("l", "")})

        chart_drawn = False
        chart_source = td_points if len(td_points) >= 2 else metrics
        if len(chart_source) >= 2:
            chart_drawn = self._market_chart(
                slide, chart_source, p, x=Inches(0.5), y=Inches(1.85), cx=Inches(4.6), cy=Inches(4.0)
            )

        if not chart_drawn:
            # Fall back to a plain metric list where the chart would be —
            # keeps the slide legible even with too little data for a chart.
            for i, dp in enumerate((td_points or metrics)[:3]):
                _txt(slide, Inches(0.6), Inches(2.0 + i * 0.9), Inches(4.2), Inches(0.7),
                     f"{labels[i]}: {dp.get('value', '')}", size=16, bold=True, color=ring_colors[i])

        # CAGR — not part of template_data.market; stays sourced from
        # SlideContent.data_points[3] (slide 10's 4th data point per the prompt).
        if len(metrics) > 3:
            cagr_dp = metrics[3]
            _txt(slide, Inches(0.5), Inches(6.0), Inches(4.6), Inches(0.35),
                 "CAGR", size=9, bold=True, color=p["text_muted"], align=PP_ALIGN.CENTER)
            _txt(slide, Inches(0.5), Inches(6.3), Inches(4.6), Inches(0.6),
                 cagr_dp.get("value", ""), size=22, bold=True, color=p["accent4"], align=PP_ALIGN.CENTER)

        # ── Metric legend (right side) ────────────────────────────────────────
        legend_x = Inches(5.3)
        ly_start  = Inches(1.9)

        legend_points = td_points if td_points else metrics
        for i, dp in enumerate(legend_points[:3]):
            y = ly_start + Inches(i * 1.15)
            acc = ring_colors[i]

            card = _rect(slide, legend_x, y, Inches(7.3), Inches(1.0), p["surface"])
            card.line.color.rgb = p["border"]
            card.line.width = Pt(1)

            # Colored left bar
            _rect(slide, legend_x, y, Inches(0.07), Inches(1.0), acc)

            _txt(slide, legend_x + Inches(0.2), y + Inches(0.1), Inches(1.0), Inches(0.35),
                 labels[i], size=9, bold=True, color=p["text_muted"])
            _txt(slide, legend_x + Inches(1.3), y + Inches(0.08), Inches(2.5), Inches(0.45),
                 dp.get("value", ""), size=28, bold=True, color=acc)
            _txt(slide, legend_x + Inches(0.2), y + Inches(0.55), Inches(7.0), Inches(0.38),
                 dp.get("sublabel", ""), size=11, color=p["text_sec"])

        # ── Revenue bullets ───────────────────────────────────────────────────
        if s.bullet_points:
            bullets = s.bullet_points[:3]
            bw = Inches(4.0)
            for i, bp in enumerate(bullets):
                bx = Inches(0.5 + i * 4.25)
                bcard = _rect(slide, bx, Inches(6.3), bw, Inches(0.75), p["surface"])
                bcard.line.color.rgb = p["border"]
                bcard.line.width = Pt(1)
                _txt(slide, bx + Inches(0.15), Inches(6.4), bw - Inches(0.3), Inches(0.55),
                     bp, size=12, color=p["text_sec"], wrap=True)

    def _market_chart(self, slide, points: list, p: dict, x, y, cx, cy) -> bool:
        """Native horizontal bar chart for TAM/SAM/SOM. `points` is up to 3
        dicts shaped {value, sublabel} in TAM/SAM/SOM order — bar length comes
        from the parsed numeric value, while the original formatted string
        ("$47B") stays the visible data label."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        cats, vals, display = [], [], []
        for label, point in zip(("TAM", "SAM", "SOM"), points):
            parsed = _parse_money_to_millions(point.get("value"))
            if parsed is None:
                continue
            cats.append(label)
            vals.append(parsed)
            display.append(point.get("value", ""))
        if len(vals) < 2:
            return False

        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("Market size", vals)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, cd).chart
        chart.has_legend = False
        chart.value_axis.visible = False
        chart.category_axis.tick_labels.font.size = Pt(13)
        chart.category_axis.tick_labels.font.bold = True
        chart.category_axis.tick_labels.font.color.rgb = p["text_pri"]
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.color.rgb = p["text_pri"]
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = p["accent1"]
        for point, label in zip(plot.series[0].points, display):
            point.data_label.text_frame.text = label
        return True

    def _timeline(self, slide, s: SlideContent, p: dict):
        """Horizontal milestone timeline."""
        self._slide_header(slide, s, p)

        items = (s.bullet_points or [])[:5]
        if not items:
            return

        n = len(items)
        item_w = Inches((13.333 - 1.4) / n)
        y_line = Inches(3.8)
        y_cards = Inches(4.1)

        line = slide.shapes.add_shape(1, Inches(0.7), y_line, W - Inches(1.4), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = p["border"]
        line.line.fill.background()

        for i, item in enumerate(items):
            x = Inches(0.7 + i * item_w.inches)
            dot_x = x + item_w / 2 - Inches(0.12)
            dot = slide.shapes.add_shape(9, dot_x, y_line - Inches(0.1), Inches(0.24), Inches(0.24))
            dot.fill.solid()
            dot.fill.fore_color.rgb = p["accent1"]
            dot.line.fill.background()

            parts = item.split(":", 1)
            label = parts[0].strip()
            detail = parts[1].strip() if len(parts) > 1 else ""

            _txt(slide, x + Inches(0.1), y_cards, item_w - Inches(0.2), Inches(0.45),
                 label, size=13, bold=True, color=p["text_pri"], align=PP_ALIGN.CENTER)
            if detail:
                _txt(slide, x + Inches(0.1), y_cards + Inches(0.45), item_w - Inches(0.2), Inches(0.6),
                     detail, size=11, color=p["text_sec"], align=PP_ALIGN.CENTER)

    def _title_bullets(self, slide, s: SlideContent, p: dict):
        """Standard title + bullets layout."""
        self._slide_header(slide, s, p)

        bullets = s.bullet_points[:8]
        body = s.body

        if body:
            _txt(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.85),
                 body, size=15, color=p["text_sec"], wrap=True)

        if not bullets:
            return

        y_start = Inches(2.95) if body else Inches(2.1)
        mid = (len(bullets) + 1) // 2
        left_bullets = bullets[:mid]
        right_bullets = bullets[mid:]
        col_w = Inches(5.7) if right_bullets else Inches(11.7)

        for i, bp in enumerate(left_bullets):
            self._bullet_row(slide, Inches(0.8), y_start + Inches(i * 0.62), col_w, bp, p["accent1"], p)

        if right_bullets:
            rx = Inches(0.8) + col_w + Inches(0.5)
            for i, bp in enumerate(right_bullets):
                self._bullet_row(slide, rx, y_start + Inches(i * 0.62), col_w, bp, p["accent2"], p)

    def _bullet_row(self, slide, x, y, w, text: str, accent, p: dict):
        """Single bullet row with colored dot."""
        dot = slide.shapes.add_shape(9, x, y + Inches(0.14), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        _txt(slide, x + Inches(0.22), y, w - Inches(0.22), Inches(0.56),
             text, size=14, color=p["text_sec"], wrap=True)

    # ── Image helpers ──────────────────────────────────────────────────────────

    def _add_background_image(self, slide, img_bytes: bytes, p: dict, overlay_alpha_pct: float = 50):
        """Add a full-slide background photo with a real (not near-opaque)
        semi-transparent overlay for text readability, tinted to the theme's
        mode so light themes get a light wash instead of always-dark."""
        try:
            img_stream = io.BytesIO(img_bytes)
            pic = slide.shapes.add_picture(img_stream, 0, 0, W, H)
            # Move picture to the very back of the shape stack
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)

            overlay = slide.shapes.add_shape(1, 0, 0, W, H)
            overlay_color = p["bg"] if p.get("mode") == "light" else _rgb("#000000")
            _set_fill_alpha(overlay, overlay_color, overlay_alpha_pct)
            overlay.line.fill.background()

            # Move overlay just behind all content but above the photo
            sp_tree.remove(overlay._element)
            sp_tree.insert(3, overlay._element)
        except Exception:
            pass  # Background image is optional — never block rendering

    def _add_logo_small(self, slide, logo_bytes: bytes):
        """Add a small, aspect-correct logo to the top-right corner of a slide."""
        try:
            native = _image_size(logo_bytes)
            logo_w, logo_h = _fit_box(native, Inches(1.3), Inches(0.5)) if native else (Inches(1.1), Inches(0.45))
            slide.shapes.add_picture(
                io.BytesIO(logo_bytes),
                W - logo_w - Inches(0.2),
                Inches(0.1) + (Inches(0.5) - logo_h) / 2,
                logo_w,
                logo_h,
            )
        except Exception:
            pass

    # ── Shared header ──────────────────────────────────────────────────────────

    def _slide_header(self, slide, s: SlideContent, p: dict):
        """Left accent bar + slide type tag + title + optional subtitle."""
        _rect(slide, 0, 0, Inches(0.07), H, p["accent1"])

        tag = s.slide_type.replace("_", " ").upper()
        _txt(slide, Inches(0.8), Inches(0.18), Inches(5), Inches(0.38),
             tag, size=9, bold=True, color=p["text_muted"])

        title_text = s.title.upper() if p.get("title_upper") else s.title
        _txt(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.85),
             title_text, size=30, bold=True, color=p["text_pri"], font=p["font_title"])

        if s.subtitle:
            _txt(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 s.subtitle, size=16, color=p["text_sec"])


# ── Low-level drawing helpers ──────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _image_size(img_bytes: bytes) -> Optional[tuple]:
    """Return (width_px, height_px) for already-Pillow-decodable image bytes,
    or None if the format can't be decoded (e.g. an SVG that slipped through)."""
    try:
        from PIL import Image
        with Image.open(io.BytesIO(img_bytes)) as img:
            return img.width, img.height
    except Exception:
        return None


def _fit_box(native: tuple, max_w: Emu, max_h: Emu) -> tuple:
    """Aspect-correct fit-within-box sizing — replaces stretching a logo to a
    fixed rectangle, which distorts any non-matching aspect ratio."""
    native_w, native_h = native
    if not native_w or not native_h:
        return max_w, max_h
    scale = min(max_w / native_w, max_h / native_h)
    return Emu(int(native_w * scale)), Emu(int(native_h * scale))


def _txt(
    slide, left, top, width, height,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_PRI,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
    font: str = "Calibri",
):
    if not text:
        return None
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align

    # Handle bold lead-ins: "**Bold part:** rest of text"
    if "**" in text:
        parts = text.split("**")
        for pi, part in enumerate(parts):
            run = p.add_run()
            run.text = part
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
            run.font.bold = bold or (pi % 2 == 1)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
        run.font.bold = bold

    return tb
