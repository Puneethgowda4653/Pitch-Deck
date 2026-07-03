"""
Slide Factory — Creates individual slides for each type/layout.

Each slide type has a dedicated builder method that:
1. Adds a blank slide
2. Applies background fill
3. Places title, content, and decorative elements
4. Applies theme fonts and colors

pptx shape positioning uses EMUs (English Metric Units):
  1 inch = 914400 EMU
  Use Inches() helper for readability.
"""
from typing import Any, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from app.ppt.engine.theme_engine import ThemeEngine


class SlideFactory:
    """Creates slides for each type using the provided theme."""

    def __init__(self, prs: Presentation, theme: ThemeEngine):
        self.prs = prs
        self.theme = theme
        # Use blank slide layout
        self.blank_layout = prs.slide_layouts[6]

    def create_slide(self, slide_data: dict) -> Any:
        """Route to the correct builder based on layout."""
        layout = slide_data.get("layout", "title_content")
        slide_type = slide_data.get("type", "custom")

        builder_map = {
            "full_bleed": self._build_full_bleed,
            "big_number": self._build_big_number,
            "title_bullets": self._build_title_bullets,
            "title_content": self._build_title_content,
            "two_column": self._build_two_column,
            "chart_focus": self._build_chart_focus,
            "image_right": self._build_title_content,  # fallback
            "image_left": self._build_title_content,   # fallback
        }

        builder = builder_map.get(layout, self._build_title_content)
        return builder(slide_data)

    # ── Slide Builders ────────────────────────────────────────────────────────

    def _build_full_bleed(self, data: dict) -> Any:
        """Cover slide — full dark background with centered title."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)
        self._add_gradient_bar(slide)

        W, H = Inches(13.33), Inches(7.5)

        # Company name / title — large centered
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.2), Inches(11.33), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data.get("title", "")
        run.font.bold = True
        run.font.size = self.theme.title_font_size("full_bleed")
        run.font.color.rgb = self.theme.text_primary
        run.font.name = self.theme.font_title

        # Subtitle / tagline
        if data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.3), Inches(9.33), Inches(1)
            )
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = data["subtitle"]
            run2.font.size = Pt(20)
            run2.font.color.rgb = self.theme.text_secondary
            run2.font.name = self.theme.font_body

        self._add_slide_number(slide, data.get("order", 0))
        return slide

    def _build_big_number(self, data: dict) -> Any:
        """Stats-focused slide with large metric callouts."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)
        self._add_accent_line(slide)

        # Title
        self._add_title(slide, data.get("title", ""), y=Inches(0.4))

        # Stats row
        stats = data.get("stats") or []
        stat_count = min(len(stats), 4)
        if stat_count > 0:
            slot_width = Inches(13.33) / stat_count
            for i, stat in enumerate(stats[:4]):
                x = slot_width * i + Inches(0.2)
                # Stat value
                val_box = slide.shapes.add_textbox(x, Inches(1.8), slot_width - Inches(0.3), Inches(1.5))
                tf = val_box.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = stat.get("value", "")
                run.font.size = self.theme.stat_value_font_size()
                run.font.bold = True
                run.font.color.rgb = self.theme.accent
                run.font.name = self.theme.font_title

                # Stat label
                lbl_box = slide.shapes.add_textbox(x, Inches(3.3), slot_width - Inches(0.3), Inches(0.5))
                tf2 = lbl_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = stat.get("label", "")
                run2.font.size = self.theme.stat_label_font_size()
                run2.font.color.rgb = self.theme.text_secondary
                run2.font.name = self.theme.font_body

        # Body text below stats
        if data.get("body_text"):
            self._add_body(slide, data["body_text"], y=Inches(4.0))

        # Bullets below body
        if data.get("bullets"):
            self._add_bullets(slide, data["bullets"], y=Inches(4.8))

        self._add_slide_number(slide, data.get("order", 0))
        return slide

    def _build_title_bullets(self, data: dict) -> Any:
        """Standard title + bullet list layout."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)
        self._add_accent_line(slide)

        self._add_title(slide, data.get("title", ""))
        if data.get("subtitle"):
            self._add_subtitle(slide, data["subtitle"])

        bullets = data.get("bullets") or []
        if bullets:
            self._add_bullets(slide, bullets, y=Inches(1.9))

        if data.get("body_text"):
            y = Inches(1.9 + len(bullets) * 0.55 + 0.2)
            self._add_body(slide, data["body_text"], y=y)

        self._add_slide_number(slide, data.get("order", 0))
        return slide

    def _build_title_content(self, data: dict) -> Any:
        """Title + mixed content layout."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)
        self._add_accent_line(slide)

        self._add_title(slide, data.get("title", ""))
        if data.get("subtitle"):
            self._add_subtitle(slide, data["subtitle"])

        if data.get("body_text"):
            self._add_body(slide, data["body_text"])

        if data.get("bullets"):
            self._add_bullets(slide, data["bullets"])

        self._add_slide_number(slide, data.get("order", 0))
        return slide

    def _build_two_column(self, data: dict) -> Any:
        """Two-column layout for comparisons."""
        return self._build_title_content(data)  # TODO: implement two-column

    def _build_chart_focus(self, data: dict) -> Any:
        """Chart-heavy slide."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)
        self._add_accent_line(slide)
        self._add_title(slide, data.get("title", ""))
        if data.get("bullets"):
            self._add_bullets(slide, data["bullets"])
        # TODO: Add chart rendering with matplotlib/pptx charts
        self._add_slide_number(slide, data.get("order", 0))
        return slide

    # ── Shared Helpers ────────────────────────────────────────────────────────

    def _apply_background(self, slide: Any) -> None:
        """Apply solid dark background."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.background

    def _add_gradient_bar(self, slide: Any) -> None:
        """Add gradient accent bar at bottom for cover slides."""
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(6.8),
            Inches(13.33), Inches(0.15),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.primary
        bar.line.fill.background()

    def _add_accent_line(self, slide: Any) -> None:
        """Add subtle top accent line."""
        line = slide.shapes.add_shape(
            1,
            Inches(0.5), Inches(0.3),
            Inches(1.2), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.accent
        line.line.fill.background()

    def _add_title(self, slide: Any, text: str, y: Emu = None) -> Any:
        if y is None:
            y = Inches(0.4)
        box = slide.shapes.add_textbox(
            Inches(0.5), y, Inches(12.33), Inches(1.2)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.bold = True
        run.font.size = self.theme.title_font_size("title_content")
        run.font.color.rgb = self.theme.text_primary
        run.font.name = self.theme.font_title
        return box

    def _add_subtitle(self, slide: Any, text: str) -> Any:
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.55), Inches(12.33), Inches(0.6)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.color.rgb = self.theme.text_secondary
        run.font.name = self.theme.font_body
        return box

    def _add_body(self, slide: Any, text: str, y: Emu = None) -> Any:
        if y is None:
            y = Inches(2.0)
        box = slide.shapes.add_textbox(
            Inches(0.5), y, Inches(12.33), Inches(2.5)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = self.theme.body_font_size()
        run.font.color.rgb = self.theme.text_secondary
        run.font.name = self.theme.font_body
        return box

    def _add_bullets(
        self, slide: Any, bullets: list[str], y: Emu = None
    ) -> Any:
        if y is None:
            y = Inches(1.9)
        box = slide.shapes.add_textbox(
            Inches(0.7), y, Inches(11.8), Inches(5.0)
        )
        tf = box.text_frame
        tf.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"• {bullet_text}"
            run.font.size = self.theme.bullet_font_size()
            run.font.color.rgb = self.theme.text_primary
            run.font.name = self.theme.font_body

        return box

    def _add_slide_number(self, slide: Any, number: int) -> None:
        """Add subtle slide number bottom-right."""
        box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(number + 1)
        run.font.size = Pt(10)
        run.font.color.rgb = self.theme.text_secondary
